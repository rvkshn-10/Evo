#!/usr/bin/env python3
"""Generate Evo PhD / mentor presentation with screenshots and pipeline diagrams."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from scripts.deck_diagrams import generate_all  # noqa: E402

OUTPUT = PROJECT_ROOT / "docs" / "Evo_PhD_Deck.pptx"
IMG = PROJECT_ROOT / "docs" / "deck_images"

BG = RGBColor(12, 18, 32)
ACCENT = RGBColor(56, 189, 248)
TEXT = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
WARN = RGBColor(251, 191, 36)
GREEN = RGBColor(52, 211, 153)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def title(slide, t: str, sub: str = "", size: int = 32):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = t
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if sub:
        s = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12), Inches(0.5))
        s.text_frame.text = sub
        s.text_frame.paragraphs[0].font.size = Pt(14)
        s.text_frame.paragraphs[0].font.color.rgb = MUTED


def bullets(slide, items: list[str], top: float = 1.55, size: int = 17):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.8), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size if not item.startswith("  ") else size - 2)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def footer(slide, text: str = "Evo · evac-evo.vercel.app"):
    f = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12), Inches(0.3))
    f.text_frame.text = text
    f.text_frame.paragraphs[0].font.size = Pt(10)
    f.text_frame.paragraphs[0].font.color.rgb = MUTED


def image(slide, path: Path, left: float, top: float, width: float):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def two_col(slide, left_items: list[str], right_items: list[str]):
    bullets(slide, left_items, top=1.5, size=16)
    box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def build() -> Path:
    generate_all()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    s = blank(prs); bg(s)
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.2))
    t.text_frame.text = "Evo"
    t.text_frame.paragraphs[0].font.size = Pt(58)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub = s.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(1.5))
    stf = sub.text_frame
    stf.text = "Evacuation intelligence · multimodal hazard fusion · hybrid ML inference"
    p2 = stf.add_paragraph(); p2.text = "evac-evo.vercel.app · github.com/rvkshn-10/Evo"
    for p in stf.paragraphs:
        p.font.size = Pt(18); p.font.color.rgb = ACCENT; p.space_after = Pt(6)

    # 2 Impact
    s = blank(prs); bg(s); title(s, "Operational impact", "How Evo helps people in emergencies")
    bullets(s, [
        "Fuses live hazard telemetry with site occupancy to estimate evacuation success and duration per campus",
        "Surfaces high-risk monitoring spots before drills or real events — density, occupancy, hazard proximity",
        "Single geospatial dashboard for dispatchers, safety officers, and researchers (not fragmented feeds)",
        "Sub-second ONNX/OpenVINO inference on edge CPU or Intel Neural Compute Stick for local deployments",
        "Longitudinal disaster history for after-action review, export to SQL/CSV for institutional analysis",
    ])

    # 3 Architecture diagram
    s = blank(prs); bg(s); title(s, "System architecture")
    image(s, IMG / "diagram-architecture.png", 0.4, 1.4, 12.4); footer(s)

    # 4 Vercel + Oracle
    s = blank(prs); bg(s); title(s, "Production deployment")
    bullets(s, [
        "Frontend: Vercel hosts static Vite build (web/dist) at evac-evo.vercel.app",
        "vercel.json rewrites /api/* → Oracle VM FastAPI :8092 (TLS at Vercel edge)",
        "Backend: Oracle Always Free ARM Ubuntu — python3 main.py via systemd",
        "CORS: CORS_ORIGINS + *.vercel.app regex in main.py",
        "History DB: SQLite locally; Neon Postgres optional via DATABASE_URL",
        "Split keeps heavy Python/ML off serverless; UI deploys on every git push",
    ])

    # 5 Data pipeline diagram
    s = blank(prs); bg(s); title(s, "Backend data pipeline")
    image(s, IMG / "diagram-data-pipeline.png", 0.4, 1.35, 12.4); footer(s)

    # 6 Dashboard screenshot
    s = blank(prs); bg(s); title(s, "Live dashboard")
    image(s, IMG / "04-dashboard-desktop.png", 0.5, 1.3, 12.2)
    footer(s, "Leaflet map · hazard layers · predictions table · run modes")

    # 7 Inputs
    s = blank(prs); bg(s); title(s, "Multimodal inputs")
    two_col(s, [
        "NOAA/NWS — CAP alerts, severity, polygons",
        "USGS — significant quakes + EEW candidates",
        "GDACS — global disaster context",
        "NASA FIRMS — wildfire thermal hotspots",
        "FEMA IPAWS — public alert feed",
    ], [
        "PeopleSense — occupancy Count, Density, Volatility",
        "FCUSD reference corpus — 4,082 labeled evacuation rows",
        "Monitoring spots — GPS, category, GGV2 match rules",
        "Hazard feature builder — severity, magnitude, distance, depth",
        "All merged into a fixed-length feature vector per site",
    ])

    # 8 Alerts
    s = blank(prs); bg(s); title(s, "Alert ingestion & enrichment")
    bullets(s, [
        "GovernmentFeedSync polls feeds → output/feeds/last_sync.json",
        "AlertProcessor.collect_live_hazards() — single-pass per dashboard request",
        "Per-alert: PeopleSense overlay → nearest hazard to spot → EvacuationPredictor",
        "Auto-deploy: Extreme/Severe NOAA + EEW → PeopleSense event POST (if enabled)",
        "Risk band: high if evac_rate < 0.9 OR density > 0.85 OR occupancy > 1500",
        "Published: GET /api/dashboard JSON + optional disaster_history snapshot",
    ])

    # 9 Map
    s = blank(prs); bg(s); title(s, "Geospatial map (Leaflet)")
    bullets(s, [
        "Base: OpenStreetMap tiles; center from quakes or Sacramento default",
        "Heatmap layer: earthquake intensity (L.heatLayer) — toggleable",
        "Markers: cyan monitoring spots, yellow quake epicenters, red/orange hazard zones",
        "Filter matrix: earthquake, wildfire, flood, tornado, tsunami, severe_weather, monitoring",
        "hazard_category assigned server-side in AlertProcessor._hazard_category()",
    ])

    # 10 Run modes overview
    s = blank(prs); bg(s); title(s, "Run modes — overview")
    bullets(s, [
        "sync — feed refresh + k-NN predictions + history (default, no LLM cost)",
        "evo — Evo 1.2 hybrid ONNX/OpenVINO inference + history",
        "evo13 — research estimator with enriched reference + optional Evo 1.3 artifacts",
        "external_ai — sync + Gemini→OpenAI narrative summary (RAG-style briefing)",
        "broadcast — full 7-agent LangChain pipeline (does not write disaster history)",
    ], size=16)

    # 11 Sync
    s = blank(prs); bg(s); title(s, "Run mode: Sync only")
    bullets(s, [
        "POST /api/alerts/sync?mode=sync",
        "EvacuationPredictor: weighted k-NN (k=25) over reference JSON",
        "Distance metric: scenario + category + occupancy/density covariates",
        "No OpenAI/Gemini — suitable for continuous monitoring without API spend",
        "Persists snapshot to SQLite/Neon via save_dashboard_snapshot()",
    ])

    # 12 Evo 1.2
    s = blank(prs); bg(s); title(s, "Run mode: Evo 1.2 hybrid (production)")
    image(s, IMG / "diagram-evo-hybrid.png", 0.4, 1.25, 6.2)
    box = s.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.5))
    tf = box.text_frame; tf.word_wrap = True
    lines = [
        "Dual-head ensemble: MLP + LightGBM (OOF average)",
        "10 numeric + categorical features → StandardScaler + one-hot",
        "Hybrid policy: k-NN for success/risk; neural for time at Train Station hubs",
        "5-fold StratifiedGroupKFold CV; synthetic augmentation grouped by row_id",
        "Exported ONNX + OpenVINO IR from Colab training",
    ]
    tf.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph(); p.text = line; p.font.size = Pt(15); p.font.color.rgb = TEXT

    # 13 Evo 1.3
    s = blank(prs); bg(s); title(s, "Evo 1.3 — synthetic demo trained", "complete_research_demo · production promotion disabled")
    bullets(s, [
        "Preflight passed: 12 PeopleSense rows + 6 synthetic drill outcomes",
        "Selected: MLP + LightGBM OOF ensemble (mlp_lightgbm_oof_average)",
        "Success MAE 1.657% ✅ · Success R² 0.185 ❌ · Time MAE 0.873 min ✅ · Time R² 0.596 ❌",
        "ONNX parity ✅ · OpenVINO parity ❌ → research runtime uses verified ONNX",
        "Classification: DATA_CEILING · keep_evo1.2_hybrid",
        "Artifacts in models/evo1.3/ — illustrative until real drill labels replace synthetic data",
    ])

    # 14 External AI / RAG
    s = blank(prs); bg(s); title(s, "Run mode: External AI (RAG-style briefing)")
    bullets(s, [
        "After dashboard build: generate_intelligence_summary() via llm_router",
        "Retrieval context: alert headlines, quake counts, high-risk spot tally from live JSON",
        "Generation: Gemini 2.0 Flash → OpenAI mini → GPT-4o failover chain",
        "Output: 4–6 sentence operational briefing in API response (not stored in history DB)",
        "Analogous to RAG: structured telemetry retrieved first, then LLM synthesis — not fine-tuned on FCUSD",
    ])

    # 15 Broadcast
    s = blank(prs); bg(s); title(s, "Run mode: Full broadcast")
    image(s, IMG / "diagram-broadcast.png", 0.4, 1.35, 12.2)
    footer(s, "7 GPT-4o agents · HeyGen-ready script output")

    # 16 Training Colab
    s = blank(prs); bg(s); title(s, "Training pipeline (Google Colab)")
    bullets(s, [
        "model_training/evo1_2/train_evo1_2.py — hazard seed + reference rows + grouped CV",
        "Requirements: numpy, pandas, scikit-learn, torch, lightgbm, onnx, openvino",
        "Notebook: model_training/evo1_2/ (export to docs/colab/)",
        "Artifacts → models/evo1.2/: evo1.2.onnx, openvino/, metrics.json, validation_report.json",
        "Quality gates: MAE, R², baseline beats, ONNX/OpenVINO parity, p95 < 10 ms",
        "Phase B (Evo 1.3): drill-aligned PeopleSense + real_outcomes.json",
    ])

    # 17 Phase status
    s = blank(prs); bg(s); title(s, "Model validation status (Phase A)")
    bullets(s, [
        "Success MAE 1.647% (pass < 2%) · Success R² 0.184 (fail ≥ 0.25)",
        "Time MAE 0.830 min (pass) · Time R² 0.636 (fail ≥ 0.75)",
        "OpenVINO p95 0.341 ms (pass) · Classification: DATA_CEILING",
        "Interpretation: live covariates lack evacuation outcome labels",
        "Production policy: hybrid k-NN success + neural time at transit categories",
    ], size=16)

    # 18 Model modal screenshot
    s = blank(prs); bg(s); title(s, "Model panel — live inference UI")
    image(s, IMG / "02-model-modal.png", 0.5, 1.25, 7.5)
    box = s.shapes.add_textbox(Inches(8.2), Inches(1.4), Inches(4.6), Inches(5.5))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate([
        "◈ Model button opens modal",
        "Live DAG: Ingest → Encode → MLP/LGBM → Merge",
        "Animated packets = inference trace",
        "Hover edges for raw→scaled values",
        "Polls /api/evo/live-flow every 2.8s",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(14); p.font.color.rgb = TEXT

    # 19 Loss chart
    s = blank(prs); bg(s); title(s, "Cross-validation loss curves")
    bullets(s, [
        "Train loss (blue): combined MLP training error per epoch — rapid initial descent",
        "Val loss (gray): held-out fold error — plateau indicates generalization limit",
        "Divergence → overfitting; parallel tracks → data ceiling (limited labeled outcomes)",
        "Curves from metrics.json — grouped StratifiedGroupKFold, not a single holdout",
        "Gates require val MAE/R² thresholds before model promotion",
    ])

    # 20 Metrics glossary
    s = blank(prs); bg(s); title(s, "Metrics in the model modal")
    two_col(s, [
        "Success MAE — mean abs error on evac success %",
        "Success R² — explained variance (success head)",
        "Time MAE — minutes error on evacuation time",
        "Time R² — explained variance (time head)",
        "Runtime — ONNX Runtime vs OpenVINO backend",
    ], [
        "Gates passed — all promotion checks true",
        "Research preview — failed R² or parity gates",
        "Hybrid policy — which head (k-NN vs neural) drove live spot",
        "Feature importance — permutation on validation fold",
        "Live success/time — current dashboard prediction",
    ])

    # 21 History
    s = blank(prs); bg(s); title(s, "Disaster history & export")
    image(s, IMG / "05-history-modal.png", 0.5, 1.25, 7.0)
    box = s.shapes.add_textbox(Inches(7.8), Inches(1.4), Inches(5.0), Inches(5.5))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate([
        "Tables: disaster_snapshots, hazard_events, evacuation_predictions",
        "Written on sync/evo/evo13/external_ai runs",
        "Charts: high-risk spots, alerts/quakes per run",
        "Export: JSON, CSV ZIP, SQLite",
        "Neon Postgres when DATABASE_URL set",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(14); p.font.color.rgb = TEXT

    # 22 NCS stack
    s = blank(prs); bg(s); title(s, "Neural Compute Stick & OpenVINO")
    image(s, IMG / "diagram-ncs-stack.png", 0.4, 1.35, 12.2)
    footer(s, "OpenVINO = driver · NCS = optional USB hardware · see docs/NCS_SETUP.md")

    # 23 NCS local
    s = blank(prs); bg(s); title(s, "Edge inference with Intel NCS")
    bullets(s, [
        "NCS2 (Myriad X) or NCS1 (Myriad 2) via USB on local Mac",
        "OpenVINO MYRIAD plugin compiles same evo1.2.xml as cloud CPU path",
        "Dashboard Inference device: Auto / CPU / NCS1 / NCS2",
        "POST /api/evo/accelerator — hot-swap without API restart",
        "Vercel shows banner: cloud has no USB — CPU only in production",
        "docs/NCS_SETUP.md — full install guide (OpenVINO is the driver, not separate from NCS)",
    ])

    # 24 Vector / features
    s = blank(prs); bg(s); title(s, "Feature vector & encoding")
    bullets(s, [
        "Numeric: log-scaled occupancy, density, hazard severity/magnitude/distance/depth, …",
        "Categorical: site category, scenario, event_type, hazard_source → one-hot",
        "StandardScaler fit on training fold; schema in feature_schema.json",
        "Live-flow UI shows each edge as raw → scaled tensor values",
        "k-NN operates in hand-crafted distance space; MLP/LGBM in scaled ℝⁿ",
    ])

    # 25 Production path
    s = blank(prs); bg(s); title(s, "Path to production deployment")
    bullets(s, [
        "Now: research preview live at evac-evo.vercel.app with Evo 1.2 hybrid policy",
        "After Colab training completes: copy artifacts → models/evo1.3/, git push",
        "Oracle: git pull + systemctl restart emergency-api",
        "Vercel: auto-deploy frontend on push to main",
        "Optional: Neon DATABASE_URL for durable multi-region history",
        "NCS: local Mac demonstrations; cloud remains CPU for public URL",
    ])

    # 26 Colab placeholder
    s = blank(prs); bg(s); title(s, "Colab notebooks", "docs/colab/ — add exported notebooks here")
    bullets(s, [
        "Evo_1_2_Training.ipynb — grouped CV, ONNX/OpenVINO export",
        "Evo_1_3_Training.ipynb — Phase B supervised training with preflight guard",
        "After export: place notebooks in docs/colab/ for presentation appendix",
        "Training outputs: validation_report.json drives promotion_recommendation",
    ])

    # 27 Close
    s = blank(prs); bg(s)
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(1.0))
    t.text_frame.text = "Thank you"
    t.text_frame.paragraphs[0].font.size = Pt(48)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    links = s.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(12), Inches(2))
    ltf = links.text_frame
    for line in ["evac-evo.vercel.app", "github.com/rvkshn-10/Evo", "docs/NCS_SETUP.md · docs/Evo_PhD_Deck.pptx"]:
        p = ltf.paragraphs[0] if line == "evac-evo.vercel.app" else ltf.add_paragraph()
        p.text = line; p.font.size = Pt(20); p.font.color.rgb = ACCENT; p.space_after = Pt(8)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
