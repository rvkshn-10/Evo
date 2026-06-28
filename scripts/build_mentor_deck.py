#!/usr/bin/env python3
"""Generate Evo mentor presentation (PowerPoint)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT = PROJECT_ROOT / "docs" / "Evo_Mentor_Deck.pptx"

# Dashboard-inspired palette
BG = RGBColor(12, 18, 32)
ACCENT = RGBColor(56, 189, 248)
TEXT = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
WARN = RGBColor(251, 191, 36)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_block(
    slide,
    title: str,
    subtitle: str = "",
    *,
    title_size: int = 36,
) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.8), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(11.5), Inches(0.6))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = Pt(16)
        stf.paragraphs[0].font.color.rgb = MUTED


def add_bullets(slide, items: list[str], *, top: float = 1.9, accent_first: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.85), Inches(top), Inches(11.3), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20 if not item.startswith("  ") else 18)
        p.font.color.rgb = ACCENT if accent_first and i == 0 else TEXT
        p.space_after = Pt(10)


def add_footer(slide, text: str = "Evo · FCUSD research preview") -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(11.5), Inches(0.35))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED


def add_metric_table(slide, rows: list[tuple[str, str, str, str]], top: float = 2.0) -> None:
    table = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.9), Inches(top), Inches(11.0), Inches(0.45 * (len(rows) + 1))).table
    headers = ("Metric", "Result", "Gate", "Status")
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = ACCENT
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = TEXT if c < 3 else (ACCENT if value == "Pass" else WARN)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Title
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    title.text_frame.text = "Evo"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.0), Inches(1.0))
    sub.text_frame.text = "Evacuation intelligence for FCUSD"
    sub.text_frame.paragraphs[0].font.size = Pt(24)
    sub.text_frame.paragraphs[0].font.color.rgb = ACCENT
    meta = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.0), Inches(1.2))
    mtf = meta.text_frame
    mtf.text = "Live dashboard: evac-evo.vercel.app\nGitHub: github.com/rvkshn-10/Evo"
    for p in mtf.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED

    # 2 — Problem
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "The problem", "School emergencies need faster, data-driven decisions")
    add_bullets(
        slide,
        [
            "Hazards arrive from many feeds — weather, earthquakes, wildfire, transit",
            "Occupancy changes minute by minute across campuses",
            "Evacuation success and time depend on site layout, crowd density, and routes",
            "Staff need one view: what's happening, who's at risk, how long evacuation may take",
        ],
    )
    add_footer(slide)

    # 3 — Solution
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "What Evo does")
    add_bullets(
        slide,
        [
            "Ingests live hazard feeds (NOAA, USGS, GDACS, NASA FIRMS, FEMA IPAWS)",
            "Reads PeopleSense occupancy for FCUSD monitoring spots",
            "Predicts evacuation success % and time per site using ML (Evo 1.2 hybrid)",
            "Publishes an interactive map dashboard + run history and exports",
            "Optional AI agent pipeline for research and broadcast (Gemini / OpenAI)",
        ],
    )
    add_footer(slide)

    # 4 — Architecture
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Architecture")
    arch = slide.shapes.add_textbox(Inches(0.9), Inches(1.85), Inches(11.2), Inches(4.5))
    atf = arch.text_frame
    atf.word_wrap = True
    lines = [
        "Browser  →  Vercel (Evo dashboard UI)",
        "         →  Oracle Cloud VM :8092 (FastAPI)",
        "                ├── Hazard feeds + PeopleSense GET API",
        "                ├── Evo 1.2 ONNX / OpenVINO inference (CPU in cloud)",
        "                └── Disaster history (SQLite or Neon Postgres)",
        "",
        "Local Mac (optional): USB Neural Compute Stick via OpenVINO MYRIAD",
        "         →  localhost:5173 + localhost:8092 only",
    ]
    atf.text = lines[0]
    for line in lines[1:]:
        p = atf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT if line else MUTED
        p.space_after = Pt(4)
    atf.paragraphs[0].font.size = Pt(18)
    atf.paragraphs[0].font.color.rgb = ACCENT
    add_footer(slide)

    # 5 — Dashboard
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Live dashboard", "https://evac-evo.vercel.app")
    add_bullets(
        slide,
        [
            "Map layers: earthquakes, weather/hazards, monitoring spots, heatmap",
            "Evacuation predictions table: occupancy, evac rate, time, risk, model",
            "Run modes: Sync (free) · Evo 1.2 hybrid (production) · Evo 1.3 (research)",
            "Model panel: live inference trace, validation metrics, loss curves",
            "History: charts, high-risk log, CSV/ZIP export",
        ],
    )
    add_footer(slide)

    # 6 — Data sources
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Data sources")
    add_bullets(
        slide,
        [
            "NOAA / NWS — active alerts for California",
            "USGS — earthquake early warning feed",
            "GDACS · NASA FIRMS · FEMA IPAWS — global hazard context",
            "PeopleSense — FCUSD occupancy (Raspberry Pi → PeopleSense DB → our GET API)",
            "FCUSD reference dataset — 4,082 evacuation modeling records (drill history)",
        ],
    )
    add_footer(slide)

    # 7 — Evo 1.2
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Evo 1.2 hybrid (production policy)")
    add_bullets(
        slide,
        [
            "Dual-head model: evacuation success % + evacuation time (minutes)",
            "Hybrid policy: k-NN for success where data is thin; neural net for time on Train Station",
            "Trained on FCUSD reference + live hazard features — no paid API per refresh",
            "Exported ONNX + OpenVINO for fast CPU inference (<1 ms p95 on cloud)",
            "Honest limitation: success R² hits a data ceiling without measured drill outcomes",
        ],
    )
    add_footer(slide)

    # 8 — Metrics
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Evo 1.2 validation (Phase A)", "Refreshed June 2026")
    add_metric_table(
        slide,
        [
            ("Success MAE", "1.647%", "< 2%", "Pass"),
            ("Success R²", "0.184", "≥ 0.25", "Fail"),
            ("Time MAE", "0.830 min", "< 1 min", "Pass"),
            ("Time R²", "0.636", "≥ 0.75", "Fail"),
            ("OpenVINO p95", "0.341 ms", "< 10 ms", "Pass"),
        ],
    )
    note = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.0), Inches(1.2))
    note.text_frame.text = (
        "Classification: DATA_CEILING — live feeds add features, not evacuation labels.\n"
        "Production policy unchanged: Evo 1.2 hybrid on the live site."
    )
    for p in note.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED
    add_footer(slide)

    # 9 — Evo 1.3
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Evo 1.3 (research path)")
    add_bullets(
        slide,
        [
            "Supervised by measured drill outcomes — not just occupancy or hazard feeds",
            "Joins timestamped PeopleSense samples to real success % and evac time",
            "Adds egress geometry (exits, route length, blockage fraction)",
            "Trainer refuses to run until preflight passes (data-first guardrails)",
            "Dashboard mode labeled research only — not for official FCUSD operations yet",
        ],
    )
    add_footer(slide)

    # 10 — Phase B
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Phase B — what unlocks better models")
    add_bullets(
        slide,
        [
            "Measured real_outcomes.json from FCUSD evacuation drills",
            "Drill-timestamp PeopleSense XML for Vista del Lago, Folsom High, Cordova Park",
            "Confirmed GPS + GGV2 Group IDs in monitoring_locations.json",
            "Egress measurements where available",
            "Repeated drills across occupancy levels and blocked-route scenarios",
        ],
    )
    add_footer(slide)

    # 11 — Synthetic demo
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Mentor demo (synthetic Phase B)", "Pipeline proof while real drills are in progress")
    add_bullets(
        slide,
        [
            "Synthetic drill outcomes + aligned PeopleSense XML in data/incoming/evo1.3/",
            "Preflight can pass → Evo 1.3 trainer runs end-to-end",
            "Clearly labeled NOT FOR PRODUCTION — illustrative numbers only",
            "Shows mentors: the system works when FCUSD supplies real labeled data",
            "Validation gates may still fail until enough real drill rows exist",
        ],
    )
    add_footer(slide)

    # 12 — NCS
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Neural Compute Stick (local only)")
    add_bullets(
        slide,
        [
            "Intel NCS1/NCS2 USB stick works only on the Mac running python3 main.py",
            "Vercel / Oracle cloud has no USB — production inference is CPU-only",
            "OpenVINO MYRIAD plugin drives the stick on localhost:5173",
            "Cloud site shows a banner explaining this limitation",
        ],
    )
    add_footer(slide)

    # 13 — Deployment
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "Deployment today")
    add_bullets(
        slide,
        [
            "UI: Vercel — evac-evo.vercel.app (auto-deploy from GitHub)",
            "API: Oracle Always Free ARM VM (FastAPI on port 8092)",
            "Default run mode: Sync only (no API credits)",
            "Recommended for demos: Evo 1.2 hybrid (production policy)",
            "Research preview disclaimer in dashboard footer",
        ],
    )
    add_footer(slide)

    # 14 — Ask
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "What we need from FCUSD", "To move from research preview → validated operations")
    add_bullets(
        slide,
        [
            "Schedule evacuation drills with measured outcomes at key sites",
            "PeopleSense exports at drill start, during, and completion",
            "Confirm site coordinates and GGV2 sensor mapping",
            "Review dashboard with safety staff — feedback on predictions and UX",
            "Operational sign-off before any live evacuation decision use",
        ],
        accent_first=True,
    )
    add_footer(slide)

    # 15 — Close
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.2))
    title.text_frame.text = "Thank you"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    links = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.0), Inches(2.0))
    ltf = links.text_frame
    ltf.text = "Dashboard: evac-evo.vercel.app"
    p2 = ltf.add_paragraph()
    p2.text = "Code: github.com/rvkshn-10/Evo"
    p3 = ltf.add_paragraph()
    p3.text = "Contact: ravipothanngaari@gmail.com"
    for p in ltf.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT
        p.space_after = Pt(8)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
